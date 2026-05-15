"""Extract plain text from common document formats for ingestion and search."""

from __future__ import annotations

import logging
import re
import shutil
import subprocess
import tempfile
from io import BytesIO
from pathlib import Path

from ai_search_assistant.ingestion.document_formats import (
    SUPPORTED_EXTENSIONS,
    TEXT_EXTENSIONS,
    extension_of,
    is_supported_filename,
    supported_extensions,
)

logger = logging.getLogger(__name__)

_WHITESPACE = re.compile(r"\n{3,}")


def normalize_extracted_text(text: str) -> str:
    """Collapse excessive blank lines; strip ends."""
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = _WHITESPACE.sub("\n\n", text)
    return text.strip()


def extract_text_from_bytes(data: bytes, filename: str) -> str:
    """Parse file bytes to UTF-8 plain text using the file extension."""
    ext = extension_of(filename)
    if ext not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValueError(
            f"Unsupported file type {ext or '(none)'!r} for {filename!r}. "
            f"Supported: {supported}"
        )
    if not data:
        raise ValueError(f"File is empty: {filename!r}")

    if ext in TEXT_EXTENSIONS:
        text = data.decode("utf-8", errors="replace")
    elif ext == ".pdf":
        text = _extract_pdf(data)
    elif ext == ".docx":
        text = _extract_docx(data)
    elif ext == ".doc":
        text = _extract_doc_legacy(data, filename)
    elif ext == ".rtf":
        text = _extract_rtf(data)
    elif ext in {".html", ".htm"}:
        text = _extract_html(data)
    elif ext == ".pptx":
        text = _extract_pptx(data)
    else:
        raise ValueError(f"No extractor registered for {ext}")

    text = normalize_extracted_text(text)
    if not text:
        raise ValueError(f"No extractable text in {filename!r}")
    return text


def extract_text_from_path(path: Path) -> str:
    """Read a file from disk and extract text (used by manifest loader)."""
    return extract_text_from_bytes(path.read_bytes(), path.name)


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(data))
    parts: list[str] = []
    for page in reader.pages:
        part = page.extract_text() or ""
        if part.strip():
            parts.append(part)
    return "\n\n".join(parts)


def _extract_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(BytesIO(data))
    parts: list[str] = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def _extract_rtf(data: bytes) -> str:
    from striprtf.striprtf import rtf_to_text

    return rtf_to_text(data.decode("utf-8", errors="replace"))


def _extract_html(data: bytes) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(data, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    return soup.get_text(separator="\n")


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(BytesIO(data))
    parts: list[str] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_bits: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_bits.append(shape.text.strip())
        if slide_bits:
            parts.append(f"Slide {slide_num}:\n" + "\n".join(slide_bits))
    return "\n\n".join(parts)


def _extract_doc_legacy(data: bytes, filename: str) -> str:
    """Legacy Word ``.doc`` via LibreOffice or antiword when installed on the host."""
    if shutil.which("soffice"):
        return _convert_via_soffice(data, filename)
    if shutil.which("antiword"):
        return _convert_via_antiword(data)
    raise ValueError(
        f"Cannot read legacy Word file {filename!r} (.doc). "
        "Install LibreOffice (soffice) in the container/host, or save as .docx."
    )


def _convert_via_soffice(data: bytes, filename: str) -> str:
    with tempfile.TemporaryDirectory() as tmp:
        root = Path(tmp)
        src = root / Path(filename).name
        src.write_bytes(data)
        out_dir = root / "out"
        out_dir.mkdir()
        cmd = [
            "soffice",
            "--headless",
            "--norestore",
            "--convert-to",
            "txt:Text",
            "--outdir",
            str(out_dir),
            str(src),
        ]
        proc = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=120,
            check=False,
        )
        if proc.returncode != 0:
            logger.warning("soffice stderr: %s", proc.stderr)
            raise ValueError(f"LibreOffice failed to convert {filename!r}")
        txt_files = list(out_dir.glob("*.txt"))
        if not txt_files:
            raise ValueError(f"LibreOffice produced no output for {filename!r}")
        return txt_files[0].read_text(encoding="utf-8", errors="replace")


def _convert_via_antiword(data: bytes) -> str:
    with tempfile.NamedTemporaryFile(suffix=".doc", delete=True) as tmp:
        tmp.write(data)
        tmp.flush()
        proc = subprocess.run(
            ["antiword", tmp.name],
            capture_output=True,
            text=True,
            timeout=60,
            check=False,
        )
        if proc.returncode != 0:
            raise ValueError("antiword failed to extract .doc text")
        return proc.stdout or ""
